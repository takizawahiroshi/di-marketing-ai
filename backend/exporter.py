"""
リッチ export — DOCX / PPTX / PDF 生成。
python-docx / python-pptx / reportlab を使用。
"""
from __future__ import annotations
import io
import re


# ── ヘルパー ──────────────────────────────────────────────────
def _agent_label(task_data: dict, agent_id: str) -> str:
    names = task_data.get("agent_names", {})
    return names.get(agent_id, agent_id)


def _iter_lines(text: str):
    """Markdown テキストを行ごとに (kind, content) で yield する。
    kind: 'h1' / 'h2' / 'h3' / 'bullet' / 'blank' / 'body'
    """
    for raw in text.split("\n"):
        line = raw.strip()
        if not line:
            yield ("blank", "")
        elif line.startswith("### "):
            yield ("h3", line[4:])
        elif line.startswith("## "):
            yield ("h2", line[3:])
        elif line.startswith("# "):
            yield ("h1", line[2:])
        elif line.startswith("- ") or line.startswith("* "):
            yield ("bullet", line[2:])
        else:
            # インラインマークダウン（**bold** など）を除去して plain text に
            plain = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
            plain = re.sub(r"\*(.+?)\*", r"\1", plain)
            plain = re.sub(r"`(.+?)`", r"\1", plain)
            yield ("body", plain)


# ── DOCX ──────────────────────────────────────────────────────
def export_to_docx(task_data: dict) -> bytes:
    from docx import Document
    from docx.shared import RGBColor, Pt

    doc = Document()

    # スタイル調整: 標準フォントを日本語対応に
    for style in doc.styles:
        if hasattr(style, "font") and style.font:
            try:
                style.font.name = "Yu Gothic UI"
            except Exception:
                pass

    # タイトル
    goal = task_data.get("goal", "DI Marketing AI Output")
    title = doc.add_heading(goal, level=0)
    if title.runs:
        title.runs[0].font.color.rgb = RGBColor(0x00, 0x17, 0xC1)

    # エージェント一覧
    agents = task_data.get("agents", [])
    if agents:
        p = doc.add_paragraph()
        run = p.add_run("使用エージェント: ")
        run.bold = True
        p.add_run(" → ".join(_agent_label(task_data, a) for a in agents))

    doc.add_paragraph()

    def _add_text_block(text: str):
        for kind, content in _iter_lines(text):
            if kind == "blank":
                doc.add_paragraph()
            elif kind == "h1":
                doc.add_heading(content, level=2)
            elif kind in ("h2", "h3"):
                doc.add_heading(content, level=3)
            elif kind == "bullet":
                doc.add_paragraph(content, style="List Bullet")
            else:
                doc.add_paragraph(content)

    # 各エージェント出力
    results = task_data.get("results", [])
    if results:
        doc.add_heading("エージェント別アウトプット", level=1)
        for r in results:
            label = _agent_label(task_data, r.get("agent_id", ""))
            doc.add_heading(label, level=2)
            _add_text_block(r.get("output", ""))

    # 統合アウトプット
    synthesis = task_data.get("synthesis", "")
    if synthesis:
        doc.add_heading("統合アウトプット", level=1)
        _add_text_block(synthesis)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX ──────────────────────────────────────────────────────
def export_to_pptx(task_data: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BLUE  = RGBColor(0x00, 0x17, 0xC1)
    GOLD  = RGBColor(0xE8, 0xB8, 0x4B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK  = RGBColor(0x1A, 0x1D, 0x23)
    MUTED = RGBColor(0x55, 0x55, 0x77)

    blank = prs.slide_layouts[6]  # 完全ブランク

    def _set_para(tf, text: str, size: int, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        p = tf.add_paragraph() if len(tf.paragraphs) > 0 and tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        p.alignment = align
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color

    def _add_title_slide(goal: str, subtitle: str):
        slide = prs.slides.add_slide(blank)
        # 背景を DI BLUE に
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BLUE

        # タイトル
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = goal
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            p.runs[0].font.size = Pt(28)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = WHITE

        # サブタイトル（エージェント列）
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(0.8))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        if p2.runs:
            p2.runs[0].font.size = Pt(13)
            p2.runs[0].font.color.rgb = GOLD

        # DI ロゴテキスト
        tb3 = slide.shapes.add_textbox(Inches(0.4), Inches(6.8), Inches(2), Inches(0.4))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "DI MARKETING AI"
        if p3.runs:
            p3.runs[0].font.size = Pt(9)
            p3.runs[0].font.color.rgb = RGBColor(0xAA, 0xBB, 0xFF)

    def _add_content_slide(heading: str, body_text: str, is_synthesis=False):
        slide = prs.slides.add_slide(blank)

        # アクセントバー
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0),
            Inches(0.18), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD if is_synthesis else BLUE
        bar.line.fill.background()

        # 見出し
        th = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.9))
        tf_h = th.text_frame
        p = tf_h.paragraphs[0]
        p.text = heading
        if p.runs:
            p.runs[0].font.size = Pt(22)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = GOLD if is_synthesis else BLUE

        # 区切り線（シミュレート用スペース）
        # 本文
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.75))
        tf = tb.text_frame
        tf.word_wrap = True

        lines = body_text.split("\n")
        first = True
        for raw in lines[:32]:  # 32行に制限
            line = raw.strip()
            if not line:
                if not first:
                    p_blank = tf.add_paragraph()
                    p_blank.space_before = Pt(2)
                continue
            p_new = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False

            if line.startswith("### ") or line.startswith("## ") or line.startswith("# "):
                clean = re.sub(r"^#+\s*", "", line)
                p_new.text = clean
                if p_new.runs:
                    p_new.runs[0].font.size = Pt(13)
                    p_new.runs[0].font.bold = True
                    p_new.runs[0].font.color.rgb = DARK
            elif line.startswith("- ") or line.startswith("* "):
                clean = re.sub(r"\*\*(.+?)\*\*", r"\1", line[2:])
                p_new.text = f"• {clean}"
                if p_new.runs:
                    p_new.runs[0].font.size = Pt(11)
                    p_new.runs[0].font.color.rgb = DARK
                p_new.level = 1
            else:
                clean = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
                p_new.text = clean
                if p_new.runs:
                    p_new.runs[0].font.size = Pt(11)
                    p_new.runs[0].font.color.rgb = DARK

    # ── タイトルスライド ──────────────────────────────────────
    goal = task_data.get("goal", "DI Marketing AI Output")
    agents = task_data.get("agents", [])
    subtitle = " → ".join(_agent_label(task_data, a) for a in agents)
    _add_title_slide(goal, subtitle)

    # ── エージェントスライド ──────────────────────────────────
    results = task_data.get("results", [])
    for r in results:
        label = _agent_label(task_data, r.get("agent_id", ""))
        output = r.get("output", "")
        # 長い場合は 2 スライドに分割（最大 1200 文字/スライド）
        _CHUNK = 1200
        chunks = [output[i:i+_CHUNK] for i in range(0, min(len(output), _CHUNK * 2), _CHUNK)]
        for j, chunk in enumerate(chunks):
            heading = label if j == 0 else f"{label}（続き）"
            _add_content_slide(heading, chunk)

    # ── 統合アウトプットスライド ──────────────────────────────
    synthesis = task_data.get("synthesis", "")
    if synthesis:
        _CHUNK = 1200
        chunks = [synthesis[i:i+_CHUNK] for i in range(0, min(len(synthesis), _CHUNK * 3), _CHUNK)]
        for j, chunk in enumerate(chunks):
            heading = "📋 統合アウトプット" if j == 0 else "📋 統合アウトプット（続き）"
            _add_content_slide(heading, chunk, is_synthesis=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF ───────────────────────────────────────────────────────
def export_to_pdf(task_data: dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    # 日本語 CID フォント（reportlab 標準添付）
    pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
    JA = "HeiseiKakuGo-W5"

    BLUE  = colors.HexColor("#0017C1")
    GOLD  = colors.HexColor("#E8B84B")
    DARK  = colors.HexColor("#1A1D23")
    MUTED = colors.HexColor("#666688")

    def _s(name, base=None, **kw) -> ParagraphStyle:
        return ParagraphStyle(name, fontName=JA, **({"parent": base} if base else {}), **kw)

    title_s = _s("T",  fontSize=20, textColor=BLUE,  leading=28, spaceAfter=8)
    sub_s   = _s("Su", fontSize=9,  textColor=MUTED, leading=13, spaceAfter=4)
    h1_s    = _s("H1", fontSize=15, textColor=BLUE,  leading=22, spaceBefore=16, spaceAfter=6)
    h2_s    = _s("H2", fontSize=12, textColor=DARK,  leading=18, spaceBefore=10, spaceAfter=4)
    h3_s    = _s("H3", fontSize=10, textColor=DARK,  leading=15, spaceBefore=6,  spaceAfter=2)
    body_s  = _s("B",  fontSize=9,  textColor=DARK,  leading=14, spaceAfter=3)
    bull_s  = _s("BL", fontSize=9,  textColor=DARK,  leading=14, spaceAfter=3,
                 leftIndent=12, firstLineIndent=0)

    def _safe(text: str) -> str:
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def _build_story(text: str, story: list):
        for kind, content in _iter_lines(text):
            if kind == "blank":
                story.append(Spacer(1, 4))
            elif kind == "h1":
                story.append(Paragraph(_safe(content), h2_s))
            elif kind == "h2":
                story.append(Paragraph(_safe(content), h2_s))
            elif kind == "h3":
                story.append(Paragraph(_safe(content), h3_s))
            elif kind == "bullet":
                story.append(Paragraph(f"• {_safe(content)}", bull_s))
            else:
                story.append(Paragraph(_safe(content), body_s))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        rightMargin=52, leftMargin=52,
        topMargin=60,   bottomMargin=52,
        title=task_data.get("goal", "DI Marketing AI"),
        author="DI Marketing AI",
    )

    story: list = []

    goal = task_data.get("goal", "DI Marketing AI Output")
    story.append(Paragraph(_safe(goal), title_s))

    agents = task_data.get("agents", [])
    if agents:
        labels = " → ".join(_agent_label(task_data, a) for a in agents)
        story.append(Paragraph(f"エージェント: {_safe(labels)}", sub_s))

    story.append(HRFlowable(width="100%", thickness=2, color=BLUE, spaceAfter=14, spaceBefore=6))

    results = task_data.get("results", [])
    if results:
        story.append(Paragraph("エージェント別アウトプット", h1_s))
        for r in results:
            label = _agent_label(task_data, r.get("agent_id", ""))
            story.append(Paragraph(_safe(label), h2_s))
            _build_story(r.get("output", ""), story)
            story.append(Spacer(1, 8))

    synthesis = task_data.get("synthesis", "")
    if synthesis:
        story.append(HRFlowable(width="100%", thickness=1, color=GOLD,
                                spaceAfter=8, spaceBefore=16))
        story.append(Paragraph("統合アウトプット", h1_s))
        _build_story(synthesis, story)

    doc.build(story)
    return buf.getvalue()
